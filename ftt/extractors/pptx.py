from __future__ import annotations

from pathlib import Path
from typing import List, Set

from ftt.extractors.base import ExtractedContent, ImageRef
from ftt.logging_utils import FileLogger
from ftt.render import convert_office_to_pdf, render_pdf_pages


def _office_enabled(mode: str) -> bool:
    return mode.lower() != "false"


def extract_pptx(
    path: Path,
    visual_mode: str,
    render_office_mode: str,
    render_dpi: int,
    render_max_pages: int,
    visuals_dir: Path,
    work_dir: Path,
    logger: FileLogger,
) -> ExtractedContent:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX extraction") from exc

    content = ExtractedContent()
    presentation = Presentation(str(path))

    slides_with_visuals: Set[int] = set()
    image_count = 0

    for index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    content.text_parts.append(f"[Slide {index}] {text}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image = shape.image
                suffix = f".{image.ext}" if image.ext else ".bin"
                image_path = visuals_dir / f"pptx_image_{image_count:04d}{suffix}"
                image_path.write_bytes(image.blob)
                content.images.append(ImageRef(path=image_path, label=f"slide {index} image {image_count}", source="pptx"))
                slides_with_visuals.add(index)
            if getattr(shape, "has_chart", False):
                slides_with_visuals.add(index)

    if visual_mode != "embedded" and _office_enabled(render_office_mode):
        try:
            if visual_mode == "full":
                slides_to_render = list(range(1, len(presentation.slides) + 1))
            else:
                slides_to_render = sorted(slides_with_visuals)
            if slides_to_render:
                slides_to_render = slides_to_render[:render_max_pages]
                pdf_path = convert_office_to_pdf(path, work_dir, logger)
                logger.info(f"Rendering {len(slides_to_render)} PPTX slides")
                image_paths = render_pdf_pages(pdf_path, visuals_dir, slides_to_render, render_dpi, logger)
                for image_path in image_paths:
                    page_num = int(image_path.stem.split("_")[-1])
                    content.images.append(ImageRef(path=image_path, label=f"slide {page_num}", source="pptx"))
        except Exception as exc:  # noqa: BLE001
            logger.warning(f"PPTX visual rendering skipped: {exc}")

    return content
